from pptx import Presentation

prs = Presentation(r"d:\Riyanshi\01_coding\projects\37 VoiceIntro AI\template.pptx")

print(f"Total Slide Layouts in template.pptx: {len(prs.slide_layouts)}")
for idx, layout in enumerate(prs.slide_layouts):
    print(f"Layout [{idx}]: {repr(layout.name)}")
    for s in layout.shapes:
        if s.has_text_frame:
            print(f"   Shape: {repr(s.name)}, text={repr(s.text_frame.text)}")

print("\nExisting Slides in template.pptx:")
for idx, slide in enumerate(prs.slides):
    print(f"Slide [{idx+1}]: layout={repr(slide.slide_layout.name)}")
