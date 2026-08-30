from pptx import Presentation

def analyze_user_pptx():
    path = r"d:\Riyanshi\01_coding\projects\37 VoiceIntro AI\VocalisAI_EchoSphere2026 present.pptx"
    prs = Presentation(path)
    print(f"Total Slides in user's PPTX: {len(prs.slides)}")
    
    for idx, slide in enumerate(prs.slides):
        print(f"\n==================== SLIDE {idx + 1} ====================")
        print(f"Layout Name: {repr(slide.slide_layout.name)}")
        has_text = False
        for shape_idx, shape in enumerate(slide.shapes):
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    has_text = True
                    print(f"   Shape[{shape_idx}] '{shape.name}': {repr(text)}")
            if shape.has_table:
                has_text = True
                print(f"   Shape[{shape_idx}] TABLE found")
        if not has_text:
            print("   [EMPTY SLIDE]")

if __name__ == "__main__":
    analyze_user_pptx()
