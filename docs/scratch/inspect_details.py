from pptx import Presentation

def inspect_all_shapes():
    prs = Presentation(r"d:\Riyanshi\01_coding\projects\37 VoiceIntro AI\template.pptx")
    
    for s_idx, slide in enumerate(prs.slides):
        print(f"\n==================== SLIDE {s_idx + 1} ====================")
        for shape_idx, shape in enumerate(slide.shapes):
            print(f"Shape [{shape_idx}] type={shape.shape_type}, name={repr(shape.name)}")
            if shape.has_text_frame:
                for p_idx, p in enumerate(shape.text_frame.paragraphs):
                    runs_text = [run.text for run in p.runs]
                    print(f"   P[{p_idx}]: {repr(p.text)} (runs: {runs_text})")

if __name__ == "__main__":
    inspect_all_shapes()
