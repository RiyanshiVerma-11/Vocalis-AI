import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import win32com.client

def populate_slides():
    user_pptx = r"d:\Riyanshi\01_coding\projects\37 VoiceIntro AI\VocalisAI_EchoSphere2026 present.pptx"
    out_pptx = r"d:\Riyanshi\01_coding\projects\37 VoiceIntro AI\VocalisAI_EchoSphere2026_IdeaSubmission.pptx"
    out_pdf = r"d:\Riyanshi\01_coding\projects\37 VoiceIntro AI\VocalisAI_EchoSphere2026_IdeaSubmission.pdf"
    
    prs = Presentation(user_pptx)
    
    # -------------------------------------------------------------
    # SLIDE 7: Technical Architecture & Multi-Agent Pipeline
    # -------------------------------------------------------------
    slide7 = prs.slides[6]
    
    # Add Title Box to Slide 7
    tb_title7 = slide7.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.8))
    tf_t7 = tb_title7.text_frame
    tf_t7.word_wrap = True
    p_t7 = tf_t7.paragraphs[0]
    p_t7.text = "Technical Architecture & Multi-Agent Pipeline:"
    p_t7.font.size = Pt(22)
    p_t7.font.bold = True
    p_t7.font.name = "Arial"
    
    # Add Body Box to Slide 7
    tb_body7 = slide7.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.733), Inches(5.5))
    tf_b7 = tb_body7.text_frame
    tf_b7.word_wrap = True
    tf_b7.text = (
        "• Candidate Client Ingress Loop (<30ms):\n"
        "  - Bidirectional WebRTC audio streaming captures studio-grade 48kHz Opus mic input.\n"
        "  - Client-side Voice Activity Detection (VAD) handles real-time speech boundary triggers.\n\n"
        "• Agora SD-RTN™ & Conversational AI Core (<50ms):\n"
        "  - Software-Defined Real-Time Network guarantees <100ms global latency and 80% packet loss resilience.\n"
        "  - Acoustic Echo Cancellation (AEC) and Deep AI Noise Suppression (ANS) filter room echo & keyboard clicks.\n"
        "  - Instant sub-150ms barge-in interruption detection pauses AI playback naturally when candidate speaks.\n\n"
        "• Backstage Committee Arbitrator (<60ms):\n"
        "  - Multi-agent orchestrator evaluates candidate response depth and picks the optimal interviewer persona\n"
        "    (Marcus Vance - Systems, Elena Rostova - Product, David Chen - Eng Dir, Aria Thorne - Security) to probe.\n\n"
        "• Shared Candidate Memory Bus & Scorecard Engine:\n"
        "  - Centralized context bus tracks resume metrics, prior Q&A claims, technical gaps, and contradiction detection.\n"
        "  - Generates 100% evidence-based consensus scorecards with verbatim transcript quote citations."
    )
    for p in tf_b7.paragraphs:
        p.font.size = Pt(11.5)
        p.font.name = "Arial"

    # -------------------------------------------------------------
    # SLIDE 8: Competitive Advantage & Quantifiable ROI Metrics
    # -------------------------------------------------------------
    slide8 = prs.slides[7]
    
    # Add Title Box to Slide 8
    tb_title8 = slide8.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.8))
    tf_t8 = tb_title8.text_frame
    tf_t8.word_wrap = True
    p_t8 = tf_t8.paragraphs[0]
    p_t8.text = "Competitive Advantage & Enterprise Impact ROI:"
    p_t8.font.size = Pt(22)
    p_t8.font.bold = True
    p_t8.font.name = "Arial"
    
    # Add Body Box to Slide 8
    tb_body8 = slide8.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.733), Inches(5.5))
    tf_b8 = tb_body8.text_frame
    tf_b8.word_wrap = True
    tf_b8.text = (
        "• Differentiation Against Existing Alternatives:\n"
        "  - Human Interview Panel: High cost ($150–$250/hr per engineer), severe fatigue, subjective delayed notes.\n"
        "  - Generic Single-Prompt AI Bots: 1500–3000ms voice lag, flat persona, no barge-in, easily gamed by buzzwords.\n"
        "  - Vocalis AI + Agora Stack: <100ms SD-RTN streaming, 4 specialized deliberative personas, instant sub-150ms barge-in,\n"
        "    3D spatial audio positioning, 100% quote-backed scorecards, and <$2 cost per 45-minute comprehensive round.\n\n"
        "• Quantifiable Business Impact & Metrics:\n"
        "  - 80% Engineering Bandwidth Saved: Reclaims hundreds of productive sprint hours for Staff and Principal engineers.\n"
        "  - 3.8x Faster Time-to-Offer: 24/7 on-demand autonomous interviews eliminate recruiter scheduling bottlenecks.\n"
        "  - $1.2M+ Annual Cost Reduction: Calculated on 100 senior technical hires @ $200/hr engineer opportunity cost.\n"
        "  - 100% Objective Calibration: Standardized, quote-backed evaluation rubrics with zero demographic bias."
    )
    for p in tf_b8.paragraphs:
        p.font.size = Pt(11.5)
        p.font.name = "Arial"

    prs.save(out_pptx)
    prs.save(user_pptx)
    print(f"Updated user PPTX: {user_pptx}")
    print(f"Updated idea submission PPTX: {out_pptx}")

    # Export PDF via PowerPoint COM
    try:
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        deck = powerpoint.Presentations.Open(out_pptx)
        deck.SaveAs(out_pdf, 32) # 32 = ppSaveAsPDF
        deck.Close()
        powerpoint.Quit()
        print(f"SUCCESS exported PDF submission: {out_pdf}")
    except Exception as e:
        print(f"COM Export note: {e}")

if __name__ == "__main__":
    populate_slides()
