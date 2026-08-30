import os
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches

def render_exact_html_slides():
    html_path = r"d:\Riyanshi\01_coding\projects\37 VoiceIntro AI\presentation.html"
    with open(html_path, "r", encoding="utf-8") as f:
        html_content = f.read()

    tmp_dir = r"d:\Riyanshi\01_coding\projects\37 VoiceIntro AI\scratch\slide_renders"
    os.makedirs(tmp_dir, exist_ok=True)

    edge_path = r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe"
    
    # Hide controls header, footer bar, and speaker notes drawer, and ensure exact 16:9 1920x1080 centering
    override_style = """
    <style>
      header, footer, #notesDrawer { display: none !important; }
      body {
        margin: 0 !important;
        padding: 0 !important;
        width: 1920px !important;
        height: 1080px !important;
        overflow: hidden !important;
        background-color: #020617 !important;
        display: flex !important;
        align-items: center !important;
        justify-content: center !important;
      }
      #slideDeck {
        width: 1920px !important;
        height: 1080px !important;
        max-width: 1920px !important;
        padding: 40px 80px !important;
        display: flex !important;
        align-items: center !important;
        justify-content: center !important;
      }
      .slide {
        width: 100% !important;
        max-width: 1760px !important;
      }
    </style>
    </head>
    """

    png_files = []
    
    for slide_num in range(1, 10):
        # 1. Inject override style into head
        modified_html = html_content.replace('</head>', override_style)
        
        # 2. Change `let currentSlide = 1;` to `let currentSlide = slide_num;`
        modified_html = modified_html.replace('let currentSlide = 1;', f'let currentSlide = {slide_num};')
        
        temp_html_path = os.path.join(tmp_dir, f"render_slide_{slide_num}.html")
        with open(temp_html_path, "w", encoding="utf-8") as f:
            f.write(modified_html)
            
        png_path = os.path.join(tmp_dir, f"final_slide_{slide_num}.png")
        
        cmd = [
            edge_path,
            "--headless",
            "--disable-gpu",
            "--hide-scrollbars",
            "--window-size=1920,1080",
            "--virtual-time-budget=2000",
            f"--screenshot={png_path}",
            f"file:///{temp_html_path.replace(os.sep, '/')}"
        ]
        
        print(f"Rendering exact Slide {slide_num}/9 screenshot...")
        subprocess.run(cmd, check=True)
        png_files.append(png_path)

    # Assemble into PowerPoint PPTX
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    
    for i, png in enumerate(png_files):
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(png, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        print(f"Added Slide {i+1} picture to PPTX deck")

    out1 = r"d:\Riyanshi\01_coding\projects\37 VoiceIntro AI\Vocalis_AI_Echosphere_Deck.pptx"
    prs.save(out1)
    print(f"SUCCESS saved exact PPTX: {out1}")


if __name__ == "__main__":
    render_exact_html_slides()
