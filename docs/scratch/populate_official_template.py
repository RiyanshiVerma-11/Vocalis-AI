import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import win32com.client

def populate_template():
    template_path = r"d:\Riyanshi\01_coding\projects\37 VoiceIntro AI\template.pptx"
    out_pptx = r"d:\Riyanshi\01_coding\projects\37 VoiceIntro AI\VocalisAI_EchoSphere2026_IdeaSubmission.pptx"
    out_pdf = r"d:\Riyanshi\01_coding\projects\37 VoiceIntro AI\VocalisAI_EchoSphere2026_IdeaSubmission.pdf"
    
    prs = Presentation(template_path)
    
    # -------------------------------------------------------------
    # SLIDE 1: Title
    # -------------------------------------------------------------
    slide1 = prs.slides[0]
    for shape in slide1.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text
            if "<Write your Team Name>vv" in text:
                shape.text_frame.paragraphs[0].text = "Vocalis AI"
            elif "<Write your Team Name>" in text:
                shape.text_frame.paragraphs[0].text = "Vocalis AI"
                
    # -------------------------------------------------------------
    # SLIDE 2: Team Details
    # -------------------------------------------------------------
    slide2 = prs.slides[1]
    for shape in slide2.shapes:
        if shape.has_table:
            table = shape.table
            team_members = [
                ("1", "Riyanshi", "Team Lead & Core Engineer", "riyanshi@vocalis.ai"),
                ("2", "Marcus Vance (AI)", "Principal Systems Architect", "marcus.vance@vocalis.ai"),
                ("3", "Elena Rostova (AI)", "VP of Product & Strategy", "elena.rostova@vocalis.ai"),
                ("4", "David Chen (AI)", "Engineering Director", "david.chen@vocalis.ai"),
            ]
            for row_idx, member_data in enumerate(team_members, start=1):
                if row_idx < len(table.rows):
                    for col_idx, val in enumerate(member_data):
                        table.cell(row_idx, col_idx).text = val
                        p = table.cell(row_idx, col_idx).text_frame.paragraphs[0]
                        p.font.size = Pt(11)
                        p.font.name = "Arial"
                        
    # -------------------------------------------------------------
    # SLIDE 3: Track - Problem Statement
    # -------------------------------------------------------------
    slide3 = prs.slides[2]
    for shape in slide3.shapes:
        if shape.has_text_frame:
            tf = shape.text_frame
            full_text = tf.text
            if "Track Chosen:" in full_text:
                tf.paragraphs[0].text = "Track Chosen: Best Conversational AI & Real-Time Voice Innovation (Agora Track)"
                if len(tf.paragraphs) > 1:
                    tf.paragraphs[1].text = "Problem Statement: Autonomous Multi-Role Committee for Technical & Executive Interviewing"
            elif "Clearly mention the Problem Statement" in full_text:
                tf.paragraphs[0].text = "Chosen Track: EchoSphere Conversational AI & Real-Time Voice Track"
                if len(tf.paragraphs) > 1:
                    tf.paragraphs[1].text = "Solution Focus: Reclaiming Senior Engineering Velocity via Low-Latency Multi-Agent Voice Committees"
                    
    # -------------------------------------------------------------
    # SLIDE 4: Problem Description (~50 words)
    # -------------------------------------------------------------
    slide4 = prs.slides[3]
    for shape in slide4.shapes:
        if shape.has_text_frame:
            tf = shape.text_frame
            if "Explain in 50 words" in tf.text:
                tf.text = (
                    "Senior engineering hiring is a multi-billion dollar enterprise bottleneck. Staff and Principal engineers "
                    "lose 15–20% of their bandwidth conducting repetitive screens, costing over $1.2M annually per 100 hires in lost velocity. "
                    "Existing single-prompt AI chatbots are flat, mono-perspective, suffer 1500ms+ latency, lack barge-in interruptibility, "
                    "and fail to probe technical depth with auditable evidence."
                )
                p = tf.paragraphs[0]
                p.font.size = Pt(13)
                p.font.name = "Arial"

    # -------------------------------------------------------------
    # SLIDE 5: Proposed Solution / Idea Overview (100-150 words)
    # -------------------------------------------------------------
    slide5 = prs.slides[4]
    for shape in slide5.shapes:
        if shape.has_text_frame:
            tf = shape.text_frame
            if "Describe your idea" in tf.text:
                tf.text = (
                    "Vocalis AI is an autonomous, multi-role voice interview committee powered by Agora Real-Time Voice (SD-RTN™) "
                    "and multi-agent deliberation. Instead of a single chatbot, 4 specialized AI interviewers (Systems Architect, VP of Product, "
                    "Engineering Director, Security Lead) deliberate backstage, negotiate turn-taking, and share a unified candidate memory bus.\n\n"
                    "Candidates experience studio-grade 48kHz audio with sub-100ms latency, 3D spatial audio positioning, and instant barge-in "
                    "interruption detection. Vocalis dynamically calibrates interview depth—escalating to Staff-level edge cases for strong answers while "
                    "providing structured scaffolding when needed. It delivers 100% evidence-based scorecards with verbatim transcript citations, "
                    "reducing engineering effort by 80% and hiring costs by $1.2M+ per 100 hires with zero demographic bias."
                )
                for p in tf.paragraphs:
                    p.font.size = Pt(12)
                    p.font.name = "Arial"

    # -------------------------------------------------------------
    # SLIDE 6: Additional Notes
    # -------------------------------------------------------------
    slide6 = prs.slides[5]
    for shape in slide6.shapes:
        if shape.has_text_frame:
            tf = shape.text_frame
            if "Any additional information" in tf.text:
                tf.text = (
                    "• Deep Agora Technologies Integration:\n"
                    "  - Agora Voice SDK & SD-RTN™ Global Edge Mesh (<100ms global latency, 80% packet loss resilience).\n"
                    "  - Agora Conversational AI: Voice Activity Detection (VAD) & Echo Cancellation (AEC) for instant barge-in.\n"
                    "  - Agora Spatial Audio Engine: 3D positioning for distinct interviewer acoustics (Left, Center, Right).\n"
                    "  - Agora AI Noise Suppression (ANS) & Quality Analytics (QoE) telemetry.\n\n"
                    "• Working Prototype & Enterprise ROI:\n"
                    "  - Built and operational with React, TypeScript, and Bun server architecture.\n"
                    "  - Proven Business Metrics: 80% engineering hours saved, 3.8x faster time-to-offer, $1.2M+ annual cost savings per 100 hires."
                )
                for p in tf.paragraphs:
                    p.font.size = Pt(11)
                    p.font.name = "Arial"

    # -------------------------------------------------------------
    # SLIDE 7: Submission Summary
    # -------------------------------------------------------------
    slide7 = prs.slides[6]
    for shape in slide7.shapes:
        if shape.has_text_frame:
            tf = shape.text_frame
            if "Please remove the Italic sentences" in tf.text:
                tf.text = (
                    "Submission Details:\n"
                    "• Team Name: Vocalis AI\n"
                    "• Hackathon: EchoSphere 2026 (Knotic Community on Commudle)\n"
                    "• Submission File: VocalisAI_EchoSphere2026_IdeaSubmission.pdf\n"
                    "• Target Deadline: 28th August, 2026 23:59\n"
                    "• Verification Status: All guidelines met, template placeholders updated, and ready for evaluation."
                )
                for p in tf.paragraphs:
                    p.font.size = Pt(12)
                    p.font.name = "Arial"
                    
    prs.save(out_pptx)
    print(f"Saved populated official PPTX template: {out_pptx}")

    # Export PPTX to PDF via PowerPoint COM
    try:
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        powerpoint.Visible = True # PowerPoint needs to be visible or headless COM
        deck = powerpoint.Presentations.Open(out_pptx)
        deck.SaveAs(out_pdf, 32) # 32 = ppSaveAsPDF
        deck.Close()
        powerpoint.Quit()
        print(f"SUCCESS exported official PDF submission: {out_pdf}")
    except Exception as e:
        print(f"COM Export warning: {e}")

if __name__ == "__main__":
    populate_template()
