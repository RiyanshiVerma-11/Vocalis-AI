from pptx import Presentation

def inspect_template():
    prs = Presentation(r"d:\Riyanshi\01_coding\projects\37 VoiceIntro AI\template.pptx")
    print(f"Total slides in template: {len(prs.slides)}")
    
    for i, slide in enumerate(prs.slides):
        print(f"\n--- SLIDE {i+1} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    print(f"  [Text]: {repr(text)}")
            if shape.has_table:
                print("  [Table found]")

if __name__ == "__main__":
    inspect_template()
