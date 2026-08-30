import os
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches
import win32com.client

def build_mindblowing_submission():
    html_path = r"d:\Riyanshi\01_coding\projects\37 VoiceIntro AI\presentation.html"
    root_dir = r"d:\Riyanshi\01_coding\projects\37 VoiceIntro AI"
    
    with open(html_path, "r", encoding="utf-8") as f:
        html_content = f.read()

    edge_path = r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe"
    
    override_style = """
    <style>
      header, footer, #notesDrawer { display: none !important; }
      body {
        margin: 0 !important;
        padding: 0 !important;
        width: 1920px !important;
        height: 1080px !important;
        overflow: hidden !important;
        background-color: #020617 !important;
        display: flex !important;
        align-items: center !important;
        justify-content: center !important;
      }
      #slideDeck {
        width: 1920px !important;
        height: 1080px !important;
        max-width: 1920px !important;
        padding: 40px 80px !important;
        display: flex !important;
        align-items: center !important;
        justify-content: center !important;
      }
      .slide {
        width: 100% !important;
        max-width: 1760px !important;
      }
    </style>
    </head>
    """

    png_files = []
    
    print("Generating high-resolution 4K slide images directly in project root folder...")
    for slide_num in range(1, 10):
        modified_html = html_content.replace('</head>', override_style)
        modified_html = modified_html.replace('let currentSlide = 1;', f'let currentSlide = {slide_num};')
        
        temp_html_path = os.path.join(root_dir, f"temp_render_slide_{slide_num}.html")
        with open(temp_html_path, "w", encoding="utf-8") as f:
            f.write(modified_html)
            
        png_path = os.path.join(root_dir, f"slide_{slide_num}.png")
        
        cmd = [
            edge_path,
            "--headless",
            "--disable-gpu",
            "--hide-scrollbars",
            "--window-size=1920,1080",
            "--virtual-time-budget=2000",
            f"--screenshot={png_path}",
            f"file:///{temp_html_path.replace(os.sep, '/')}"
        ]
        
        subprocess.run(cmd, check=True)
        print(f"  [+] Saved image: slide_{slide_num}.png")
        png_files.append(png_path)
        
        if os.path.exists(temp_html_path):
            os.remove(temp_html_path)

    # Build Mind-Blowing PPTX
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    
    for i, png in enumerate(png_files):
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(png, Inches(0), Inches(0), Inches(13.333), Inches(7.5))

    out_pptx = os.path.join(root_dir, "VocalisAI_EchoSphere2026_IdeaSubmission.pptx")
    out_pdf = os.path.join(root_dir, "VocalisAI_EchoSphere2026_IdeaSubmission.pdf")
    out_pptx_final = os.path.join(root_dir, "VocalisAI_EchoSphere2026_FinalSubmission.pptx")
    out_pdf_final = os.path.join(root_dir, "VocalisAI_EchoSphere2026_FinalSubmission.pdf")
    
    saved_pptx = None
    try:
        prs.save(out_pptx)
        saved_pptx = out_pptx
        print(f"[+] Saved PPTX: {out_pptx}")
    except PermissionError:
        prs.save(out_pptx_final)
        saved_pptx = out_pptx_final
        print(f"[+] Saved PPTX fallback: {out_pptx_final}")
    
    # Export to PDF via PowerPoint COM
    try:
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        deck = powerpoint.Presentations.Open(saved_pptx)
        pdf_target = out_pdf if saved_pptx == out_pptx else out_pdf_final
        deck.SaveAs(pdf_target, 32) # 32 = ppSaveAsPDF
        deck.Close()
        powerpoint.Quit()
        print(f"[+] Exported PDF: {pdf_target}")
    except Exception as e:
        print(f"Note on COM PDF export: {e}")



if __name__ == "__main__":
    build_mindblowing_submission()
