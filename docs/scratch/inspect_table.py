from pptx import Presentation

prs = Presentation(r"d:\Riyanshi\01_coding\projects\37 VoiceIntro AI\template.pptx")
slide2 = prs.slides[1]

for shape in slide2.shapes:
    if shape.has_table:
        table = shape.table
        print(f"Table dimensions: {len(table.rows)} rows x {len(table.columns)} cols")
        for r_idx, row in enumerate(table.rows):
            row_vals = [cell.text.strip().replace('\n', ' ') for cell in row.cells]
            print(f"Row {r_idx}: {row_vals}")
