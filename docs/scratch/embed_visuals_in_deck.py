import os
from pptx import Presentation
from pptx.util import Inches, Pt
import win32com.client

def embed_images_into_pptx():
    root_dir = r"d:\Riyanshi\01_coding\projects\37 VoiceIntro AI"
    user_pptx = os.path.join(root_dir, "VocalisAI_EchoSphere2026 present.pptx")
    out_pptx = os.path.join(root_dir, "VocalisAI_EchoSphere2026_IdeaSubmission.pptx")
    out_pdf = os.path.join(root_dir, "VocalisAI_EchoSphere2026_IdeaSubmission.pdf")
    
    arch_img = os.path.join(root_dir, "vocalis_architecture_diagram.jpg")
    dash_img = os.path.join(root_dir, "vocalis_committee_dashboard.jpg")
    
    prs = Presentation(user_pptx)
    
    # -------------------------------------------------------------
    # SLIDE 7: Technical Architecture & Diagram
    # -------------------------------------------------------------
    slide7 = prs.slides[6]
    # Add Image to Slide 7 (Right side or centered)
    # Slide width = 13.333, height = 7.5
    # Left = 6.8, Top = 1.3, Width = 5.8
    if os.path.exists(arch_img):
        # Resize existing text box on slide 7 to left side
        for s in slide7.shapes:
            if s.has_text_frame and "Technical Architecture" not in s.text_frame.text:
                s.left = Inches(0.8)
                s.width = Inches(5.8)
                s.top = Inches(1.3)
                s.height = Inches(5.5)
        
        slide7.shapes.add_picture(arch_img, Inches(6.8), Inches(1.5), Inches(5.8), Inches(3.26))

    # -------------------------------------------------------------
    # SLIDE 8: Committee Dashboard & ROI Metrics
    # -------------------------------------------------------------
    slide8 = prs.slides[7]
    if os.path.exists(dash_img):
        # Resize text box to left side
        for s in slide8.shapes:
            if s.has_text_frame and "Competitive Advantage" not in s.text_frame.text:
                s.left = Inches(0.8)
                s.width = Inches(5.8)
                s.top = Inches(1.3)
                s.height = Inches(5.5)
                
        slide8.shapes.add_picture(dash_img, Inches(6.8), Inches(1.5), Inches(5.8), Inches(3.26))

    saved_path = None
    try:
        prs.save(user_pptx)
        prs.save(out_pptx)
        saved_path = out_pptx
        print(f"[+] Successfully saved PPTX with embedded visual assets: {out_pptx}")
    except PermissionError:
        alt_pptx = os.path.join(root_dir, "VocalisAI_EchoSphere2026_FinalSubmission.pptx")
        prs.save(alt_pptx)
        saved_path = alt_pptx
        print(f"[+] Saved fallback PPTX: {alt_pptx}")

    # Export PDF via PowerPoint COM
    try:
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        deck = powerpoint.Presentations.Open(saved_path)
        pdf_target = out_pdf if saved_path == out_pptx else os.path.join(root_dir, "VocalisAI_EchoSphere2026_FinalSubmission.pdf")
        deck.SaveAs(pdf_target, 32)
        deck.Close()
        powerpoint.Quit()
        print(f"[+] Exported PDF with visuals: {pdf_target}")
    except Exception as e:
        print(f"COM Export note: {e}")

if __name__ == "__main__":
    embed_images_into_pptx()
