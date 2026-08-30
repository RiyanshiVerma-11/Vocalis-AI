import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_deck_from_html():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Precise Color Tokens matching presentation.html dark theme
    BG_DARK = RGBColor(2, 6, 23)          # slate-950 #020617
    CARD_BG = RGBColor(15, 23, 42)        # slate-900 #0F172A
    CARD_ALT = RGBColor(30, 41, 59)       # slate-800 #1E293B
    CARD_BORDER = RGBColor(51, 65, 85)   # slate-700 #334155
    
    AGORA_BLUE = RGBColor(9, 157, 253)   # #099DFD
    CYAN_ACCENT = RGBColor(0, 240, 255)   # #00F0FF
    PURPLE_ACCENT = RGBColor(168, 85, 247)# #A855F7
    TEXT_WHITE = RGBColor(255, 255, 255)  # #FFFFFF
    TEXT_MUTED = RGBColor(148, 163, 184) # slate-400 #94A3B8
    TEXT_SUB = RGBColor(203, 213, 225)   # slate-300 #CBD5E1
    
    EMERALD_GREEN = RGBColor(16, 185, 129)# #10B981
    ROSE_RED = RGBColor(244, 63, 94)      # #F43F5E
    AMBER_GOLD = RGBColor(245, 158, 11)   # #F59E0B
    INDIGO_BLUE = RGBColor(99, 102, 241)  # #6366F1
    
    blank_slide_layout = prs.slide_layouts[6]
    
    def set_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_DARK
        bg.line.fill.background()

    def add_slide_header(slide, tag_text, title_text, sub_text=None, tag_color=AGORA_BLUE):
        # Category Tag
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.35), Inches(11.733), Inches(0.35))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = tag_text.upper()
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = tag_color
        p_cat.font.name = "Arial"
        
        # Main Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.65), Inches(11.733), Inches(0.6))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(24)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_WHITE
        p_title.font.name = "Arial"
        
        y_next = Inches(1.25)
        if sub_text:
            sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.18), Inches(11.733), Inches(0.4))
            tf_sub = sub_box.text_frame
            tf_sub.word_wrap = True
            p_sub = tf_sub.paragraphs[0]
            p_sub.text = sub_text
            p_sub.font.size = Pt(12)
            p_sub.font.color.rgb = TEXT_SUB
            p_sub.font.name = "Arial"
            y_next = Inches(1.55)
            
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), y_next, Inches(11.733), Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = CARD_BORDER
        line.line.fill.background()
        
        return y_next + Inches(0.15)

    # =============================================================
    # SLIDE 1: Title & Executive Summary
    # =============================================================
    slide1 = prs.slides.add_slide(blank_slide_layout)
    set_bg(slide1)
    
    # Outer Glow Card
    card1 = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.6), Inches(11.733), Inches(6.3))
    card1.fill.solid()
    card1.fill.fore_color.rgb = CARD_BG
    card1.line.color.rgb = AGORA_BLUE
    card1.line.width = Pt(1.5)
    
    # Badge
    tb_tag = slide1.shapes.add_textbox(Inches(1.2), Inches(0.9), Inches(10.933), Inches(0.4))
    p_tag = tb_tag.text_frame.paragraphs[0]
    p_tag.text = "⚡ AGORA CONVERSATIONAL AI HACKATHON & PITCH DECK"
    p_tag.font.size = Pt(11)
    p_tag.font.bold = True
    p_tag.font.color.rgb = CYAN_ACCENT
    
    # Title & Subtitle
    tb_main = slide1.shapes.add_textbox(Inches(1.2), Inches(1.35), Inches(10.933), Inches(1.4))
    tf_main = tb_main.text_frame
    tf_main.word_wrap = True
    p1 = tf_main.paragraphs[0]
    p1.text = "Vocalis AI: The Autonomous"
    p1.font.size = Pt(44)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_WHITE
    
    p2 = tf_main.add_paragraph()
    p2.text = "Multi-Role Voice Interview Panel"
    p2.font.size = Pt(44)
    p2.font.bold = True
    p2.font.color.rgb = AGORA_BLUE
    
    tb_sub = slide1.shapes.add_textbox(Inches(1.2), Inches(2.85), Inches(10.933), Inches(0.7))
    tf_sub = tb_sub.text_frame
    tf_sub.word_wrap = True
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "Ultra-realistic, adaptive engineering and leadership interview committee powered by Agora Real-Time Voice SD-RTN™ and Multi-Agent Deliberation."
    p_sub.font.size = Pt(15)
    p_sub.font.color.rgb = TEXT_SUB
    
    # 4 Pillar Cards
    pillars = [
        ("01 • PROBLEM", "40+ Eng Hours Wasted", "Senior engineers lose sprint days conducting repetitive initial screening rounds.", AGORA_BLUE),
        ("02 • SOLUTION", "Autonomous Committee", "4 specialized AI interviewers (System Architect, VP, EM, Security) deliberate in real time.", CYAN_ACCENT),
        ("03 • AGORA INTEGRATION", "<100ms SD-RTN Voice", "Sub-100ms bidirectional speech streaming, barge-in interruptibility, and spatial audio.", INDIGO_BLUE),
        ("04 • BUSINESS ROI", "80% Cost Reduction", "$1.2M saved per 100 hires with 100% timestamped quote evidence and zero bias.", EMERALD_GREEN)
    ]
    
    for i, (num_label, title, desc, col) in enumerate(pillars):
        x = Inches(1.2 + i * 2.76)
        y = Inches(3.7)
        c = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.6), Inches(2.2))
        c.fill.solid()
        c.fill.fore_color.rgb = CARD_ALT
        c.line.color.rgb = col
        c.line.width = Pt(1)
        
        tb = slide1.shapes.add_textbox(x + Inches(0.12), y + Inches(0.12), Inches(2.36), Inches(1.96))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p_n = tf.paragraphs[0]
        p_n.text = num_label
        p_n.font.size = Pt(10)
        p_n.font.bold = True
        p_n.font.color.rgb = col
        
        p_t = tf.add_paragraph()
        p_t.text = title
        p_t.font.size = Pt(14)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_WHITE
        
        p_d = tf.add_paragraph()
        p_d.text = "\n" + desc
        p_d.font.size = Pt(10)
        p_d.font.color.rgb = TEXT_MUTED

    # Footer Metadata
    tb_foot = slide1.shapes.add_textbox(Inches(1.2), Inches(6.1), Inches(10.933), Inches(0.4))
    tf_f = tb_foot.text_frame
    tf_f.word_wrap = True
    p_f = tf_f.paragraphs[0]
    p_f.text = "Team: Vocalis Core Engineering   •   Target Category: Best Conversational AI & Real-Time Voice   •   Stack: Agora Voice SDK, SD-RTN™, Gemini Multi-Agent"
    p_f.font.size = Pt(11)
    p_f.font.color.rgb = TEXT_MUTED

    # =============================================================
    # SLIDE 2: Problem Statement & Industry Bottlenecks
    # =============================================================
    slide2 = prs.slides.add_slide(blank_slide_layout)
    set_bg(slide2)
    y_start = add_slide_header(
        slide2, 
        "1. Problem Statement & Industry Need", 
        "The Multi-Billion Dollar Technical Hiring Bottleneck",
        "Why traditional human interview panels don't scale and why first-generation single-prompt AI bots fail miserably.",
        ROSE_RED
    )
    
    probs = [
        ("💸 Massive Engineering Drain", 
         "Staff and Principal engineers spend 15–20% of their working hours conducting repetitive rounds, costing top tech enterprises over $1.2M annually in lost developer velocity per 100 hires.",
         "Metric: 42.5 hrs average engineering time per offer.",
         ROSE_RED),
        ("🤖 Superficial Single-Agent Bots", 
         "Existing AI mock tools act as flat, single-prompt conversationalists. They lack perspective specialization (System Arch vs Product SLA vs Leadership) and fail to probe hand-wavy claims with depth.",
         "Flaw: Easily gamed by buzzwords without deep verification.",
         AMBER_GOLD),
        ("⏱️ Unnatural Voice Latency & Rigid Turns", 
         "Legacy voice interfaces suffer from 1500ms–3000ms latency pauses, turn-taking collisions, and robotic speech. Candidates cannot naturally interrupt or correct an interviewer mid-sentence.",
         "Flaw: Lacks real-time barge-in and audio streaming sync.",
         INDIGO_BLUE)
    ]
    
    for i, (title, desc, metric, border_col) in enumerate(probs):
        x = Inches(0.8 + i * 3.98)
        y = y_start + Inches(0.1)
        w = Inches(3.78)
        h = Inches(4.3)
        
        c = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        c.fill.solid()
        c.fill.fore_color.rgb = CARD_BG
        c.line.color.rgb = border_col
        c.line.width = Pt(1.5)
        
        tb = slide2.shapes.add_textbox(x + Inches(0.2), y + Inches(0.2), w - Inches(0.4), h - Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(16)
        p1.font.bold = True
        p1.font.color.rgb = TEXT_WHITE
        
        p2 = tf.add_paragraph()
        p2.text = "\n" + desc
        p2.font.size = Pt(12)
        p2.font.color.rgb = TEXT_SUB
        
        p3 = tf.add_paragraph()
        p3.text = "\n" + metric
        p3.font.size = Pt(10)
        p3.font.bold = True
        p3.font.color.rgb = border_col

    # Problem Summary Banner
    banner = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.15), Inches(11.733), Inches(0.8))
    banner.fill.solid()
    banner.fill.fore_color.rgb = CARD_ALT
    banner.line.color.rgb = ROSE_RED
    
    tb_b = slide2.shapes.add_textbox(Inches(1.0), Inches(6.25), Inches(11.333), Inches(0.6))
    tf_b = tb_b.text_frame
    tf_b.word_wrap = True
    p_b = tf_b.paragraphs[0]
    p_b.text = "🎯 The Core Need: An ultra-low latency, multi-agent committee that actively listens, interrupts naturally via Agora, and conducts adaptive, evidence-based evaluations."
    p_b.font.size = Pt(12)
    p_b.font.bold = True
    p_b.font.color.rgb = TEXT_WHITE

    # =============================================================
    # SLIDE 3: Proposed Solution & Core Innovation
    # =============================================================
    slide3 = prs.slides.add_slide(blank_slide_layout)
    set_bg(slide3)
    y_start = add_slide_header(
        slide3,
        "2. Proposed Solution & Innovation",
        "The Autonomous Multi-Role Voice Interview Committee",
        "A next-generation platform where multiple AI specialists collaborate in real time with shared memory, dynamic follow-ups, and natural speech.",
        CYAN_ACCENT
    )
    
    sols = [
        ("01", "Autonomous Committee Turn-Taking", "Rather than one chatbot pretending to be everyone, Vocalis runs distinct specialists (Marcus Vance, Elena Rostova, David Chen). They negotiate speaking turns backstage based on who has the sharpest follow-up.", AGORA_BLUE),
        ("02", "Real-Time Shared Candidate Memory", "All agents read and write to a unified context bus containing candidate resume metrics, previous answers, unresolved technical probes, detected contradictions, and competency progression.", CYAN_ACCENT),
        ("03", "Dynamic Depth & Difficulty Calibration", "Analyzes candidate answers for superficial buzzwords vs genuine architectural depth. Strong answers trigger Staff/Principal tier edge-case escalations, while struggling candidates receive scaffolding.", INDIGO_BLUE),
        ("04", "100% Evidence-Based Scorecard Generation", "Generates post-round committee consensus scorecards with timestamped quote citations, rubric ratings, flagged architectural risks, and personalized candidate development roadmaps.", EMERALD_GREEN)
    ]
    
    card_pos_sol = [
        (Inches(0.8), y_start + Inches(0.1)),
        (Inches(6.8), y_start + Inches(0.1)),
        (Inches(0.8), y_start + Inches(2.25)),
        (Inches(6.8), y_start + Inches(2.25)),
    ]
    
    for i, (num, title, desc, col) in enumerate(sols):
        x, y = card_pos_sol[i]
        w, h = Inches(5.733), Inches(2.0)
        
        c = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        c.fill.solid()
        c.fill.fore_color.rgb = CARD_BG
        c.line.color.rgb = col
        c.line.width = Pt(1)
        
        tb = slide3.shapes.add_textbox(x + Inches(0.2), y + Inches(0.15), w - Inches(0.4), h - Inches(0.3))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p1 = tf.paragraphs[0]
        p1.text = f"{num} • {title}"
        p1.font.size = Pt(15)
        p1.font.bold = True
        p1.font.color.rgb = col
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = TEXT_SUB

    # Innovation Badges Row
    badges_tb = slide3.shapes.add_textbox(Inches(0.8), Inches(6.1), Inches(11.733), Inches(0.5))
    tf_bg = badges_tb.text_frame
    p_bg = tf_bg.paragraphs[0]
    p_bg.text = "✓ Multi-Agent Handoffs   •   ✓ Sub-100ms Interruptibility   •   ✓ Resume Verification   •   ✓ Zero Demographic Bias"
    p_bg.font.size = Pt(12)
    p_bg.font.bold = True
    p_bg.font.color.rgb = CYAN_ACCENT
    p_bg.alignment = PP_ALIGN.CENTER

    # =============================================================
    # SLIDE 4: Technical Architecture & Pipeline
    # =============================================================
    slide4 = prs.slides.add_slide(blank_slide_layout)
    set_bg(slide4)
    y_start = add_slide_header(
        slide4,
        "3. Technical Architecture & End-to-End Pipeline",
        "Full-Stack Real-Time Multi-Agent Topology",
        "Sub-120ms total roundtrip audio loop with autonomous deliberation, shared state, and instant barge-in.",
        AGORA_BLUE
    )
    
    # Architecture 4 Nodes Box
    nodes = [
        ("Candidate Client", "Browser / WebRTC\n🎤 Mic Stream In\n🔊 Multi-Voice Out\nInstant Barge-In VAD", Inches(0.8), AGORA_BLUE),
        ("Agora SD-RTN™ & SDK", "Agora Conversational AI\nVAD • AEC • ANS\nSpatial Audio Engine\nQuality Telemetry", Inches(3.8), CYAN_ACCENT),
        ("Multi-Agent Deliberation", "Marcus (Architect)\nElena (Product VP)\nDavid (Eng Dir)\nAria (Security)\nArbitration Engine", Inches(6.8), PURPLE_ACCENT),
        ("Shared Context Memory", "Resume Data Index\nMulti-Turn Q&A History\nEvidence Radar Matrix\nContradiction Tracker", Inches(9.8), EMERALD_GREEN)
    ]
    
    for title, body, x, col in nodes:
        y = y_start + Inches(0.1)
        w = Inches(2.733)
        h = Inches(3.4)
        
        c = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        c.fill.solid()
        c.fill.fore_color.rgb = CARD_BG
        c.line.color.rgb = col
        c.line.width = Pt(1.5)
        
        tb = slide4.shapes.add_textbox(x + Inches(0.15), y + Inches(0.15), w - Inches(0.3), h - Inches(0.3))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = col
        
        p2 = tf.add_paragraph()
        p2.text = "\n" + body
        p2.font.size = Pt(11)
        p2.font.color.rgb = TEXT_WHITE

    # Bottom 3 Step Timing Breakdown Cards
    timings = [
        ("1. Ingress & VAD (<30ms)", "Agora SDK captures 48kHz audio and detects speech boundaries instantly with zero buffer bloat."),
        ("2. Committee Arbitrator (<60ms)", "Multi-agent orchestrator ranks candidate depth and picks the optimal specialist to probe."),
        ("3. Egress Stream (<30ms)", "Agora SD-RTN streams selected interviewer's voice with 3D spatial positioning and barge-in.")
    ]
    
    for i, (title, desc) in enumerate(timings):
        x = Inches(0.8 + i * 3.98)
        y = Inches(5.4)
        w = Inches(3.78)
        h = Inches(1.5)
        
        c = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        c.fill.solid()
        c.fill.fore_color.rgb = CARD_ALT
        c.line.color.rgb = CARD_BORDER
        c.line.width = Pt(1)
        
        tb = slide4.shapes.add_textbox(x + Inches(0.15), y + Inches(0.15), w - Inches(0.3), h - Inches(0.3))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(12)
        p1.font.bold = True
        p1.font.color.rgb = CYAN_ACCENT
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(10)
        p2.font.color.rgb = TEXT_MUTED

    # =============================================================
    # SLIDE 5: Planned Utilization of Agora Technologies
    # =============================================================
    slide5 = prs.slides.add_slide(blank_slide_layout)
    set_bg(slide5)
    y_start = add_slide_header(
        slide5,
        "4. Deep Dive: Planned Utilization of Agora Technologies",
        "Harnessing the Full Power of Agora Real-Time AI",
        "How Vocalis natively integrates Agora's core SDKs, Conversational AI Engine, and Global SD-RTN™ mesh.",
        AGORA_BLUE
    )
    
    agora_pillars = [
        ("01", "Agora Voice SDK (Web & Mobile)", "Ultra-low latency bidirectional audio streaming with Opus 48kHz studio audio, dynamic bitrate scaling, and WebRTC fallback.", "Agora API: RtcEngine.joinChannel()", AGORA_BLUE),
        ("02", "Agora SD-RTN™ Global Mesh", "Software-Defined Real-Time Network guarantees <100ms global latency, 99.99% uptime, and up to 80% packet loss resilience.", "Feature: Dynamic Edge Routing", CYAN_ACCENT),
        ("03", "Conversational AI & Barge-In VAD", "Integrated Voice Activity Detection & Echo Cancellation. When candidate speaks, Agora immediately interrupts AI audio.", "Feature: Sub-150ms Instant Barge-In", INDIGO_BLUE),
        ("04", "Multi-Track Spatial Audio", "Assigns separate 3D spatial audio positions (Architect Left, VP Center, Security Right) replicating physical room acoustics.", "Agora API: SpatialAudioEngine", PURPLE_ACCENT),
        ("05", "Agora AI Noise Suppression (ANS)", "Deep-learning background noise suppression removes keyboard clicks, HVAC rumble, and room echo for clean STT.", "Feature: Deep AI Noise Removal", EMERALD_GREEN),
        ("06", "Agora Analytics & Quality Insights", "Real-time telemetry tracking round audio quality, candidate speech clarity, packet loss, and latency metrics to verify fairness.", "Feature: Agora QoE Telemetry", AMBER_GOLD)
    ]
    
    grid_pos = [
        (Inches(0.8), y_start + Inches(0.1)),
        (Inches(4.8), y_start + Inches(0.1)),
        (Inches(8.8), y_start + Inches(0.1)),
        (Inches(0.8), y_start + Inches(2.5)),
        (Inches(4.8), y_start + Inches(2.5)),
        (Inches(8.8), y_start + Inches(2.5)),
    ]
    
    for i, (num, title, desc, api_feat, col) in enumerate(agora_pillars):
        x, y = grid_pos[i]
        w, h = Inches(3.733), Inches(2.3)
        
        c = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        c.fill.solid()
        c.fill.fore_color.rgb = CARD_BG
        c.line.color.rgb = col
        c.line.width = Pt(1)
        
        tb = slide5.shapes.add_textbox(x + Inches(0.15), y + Inches(0.12), w - Inches(0.3), h - Inches(0.24))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p1 = tf.paragraphs[0]
        p1.text = f"{num} • {title}"
        p1.font.size = Pt(12)
        p1.font.bold = True
        p1.font.color.rgb = TEXT_WHITE
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(10)
        p2.font.color.rgb = TEXT_SUB
        
        p3 = tf.add_paragraph()
        p3.text = "\n" + api_feat
        p3.font.size = Pt(9)
        p3.font.bold = True
        p3.font.color.rgb = col

    # =============================================================
    # SLIDE 6: Competitive Advantage & Innovation Matrix
    # =============================================================
    slide6 = prs.slides.add_slide(blank_slide_layout)
    set_bg(slide6)
    y_start = add_slide_header(
        slide6,
        "5. Competitive Advantage & Innovation Matrix",
        "Why Vocalis Outperforms All Existing Alternatives",
        "Comprehensive comparison against traditional human committees and first-generation AI mock bots.",
        PURPLE_ACCENT
    )
    
    # Comparison Table
    rows, cols = 7, 4
    t_shape = slide6.shapes.add_table(rows, cols, Inches(0.8), y_start + Inches(0.1), Inches(11.733), Inches(4.9))
    table = t_shape.table
    
    table.columns[0].width = Inches(2.7)
    table.columns[1].width = Inches(2.8)
    table.columns[2].width = Inches(2.8)
    table.columns[3].width = Inches(3.433)
    
    headers_comp = ["Capability Dimension", "Human Interview Panel", "Generic AI Voice Bots", "Vocalis AI + Agora Stack"]
    for j, h in enumerate(headers_comp):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = PURPLE_ACCENT if j < 3 else AGORA_BLUE
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        
    matrix_data = [
        ("Committee Deliberation", "High, but uncoordinated & fatigued", "❌ Single flat perspective", "✓ 4 Specialized Agents + Backstage Deliberation"),
        ("Voice Latency & Flow", "Natural (0ms)", "❌ 1500–3000ms lag, unnatural", "✓ <100ms with Agora SD-RTN™ streaming"),
        ("Real-Time Barge-In", "Natural human interruption", "❌ AI talks over candidate", "✓ Sub-150ms Instant Voice Interruption (VAD)"),
        ("Resume & Metric Probing", "Often skimmed in 2 mins", "❌ Generic question list", "✓ Cites specific past claims, numbers & tech stack"),
        ("Evaluation Scorecard", "Subjective, delayed notes", "❌ Generic pass/fail summary", "✓ Timestamped quotes, radar rubrics, action plan"),
        ("Engineering Cost", "$150–$250/hr per engineer", "$10–$20/mo", "<$2 per comprehensive 45-min panel round")
    ]
    
    for i, row_vals in enumerate(matrix_data):
        for j, val in enumerate(row_vals):
            cell = table.cell(i+1, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD_BG if i % 2 == 0 else CARD_ALT
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.size = Pt(10)
            if j == 0:
                p.font.bold = True
                p.font.color.rgb = TEXT_WHITE
            elif j == 3:
                p.font.bold = True
                p.font.color.rgb = CYAN_ACCENT
            else:
                p.font.color.rgb = TEXT_MUTED

    # =============================================================
    # SLIDE 7: Technical Feasibility & Execution Roadmap
    # =============================================================
    slide7 = prs.slides.add_slide(blank_slide_layout)
    set_bg(slide7)
    y_start = add_slide_header(
        slide7,
        "6. Technical Feasibility & Execution Roadmap",
        "From Working Prototype to Enterprise Global Scale",
        "De-risked architecture with concrete deployment milestones and production-ready components.",
        EMERALD_GREEN
    )
    
    phases = [
        ("PHASE 1 • COMPLETE", "Functional Studio Prototype", ["Multi-Agent turn-taking logic", "Shared context memory & resume engine", "Real-time speech visualizer & barge-in", "Evidence scorecard generator"], EMERALD_GREEN, "100% DONE"),
        ("PHASE 2 • CURRENT", "Full Agora Conversational AI", ["Agora WebRTC Native Channel binding", "Agora Spatial Audio Persona placement", "Deep-learning Noise Suppression (ANS)", "Sub-80ms global SD-RTN edge routing"], AGORA_BLUE, "IN PROGRESS"),
        ("PHASE 3 • Q3 2026", "Enterprise ATS Integration", ["Greenhouse, Lever & Ashby 1-click sync", "Custom role rubric calibration studio", "SOC-2 Type II audit logging", "Automated panel debrief digests"], INDIGO_BLUE, "PLANNED"),
        ("PHASE 4 • Q4 2026", "Multilingual Global Mesh", ["Simultaneous multi-language interviews", "Live code execution sandbox stream", "Executive leadership board simulations", "Agora Whiteboard interactive diagramming"], PURPLE_ACCENT, "SCALE")
    ]
    
    for i, (p_tag, p_title, p_bullets, col, status) in enumerate(phases):
        x = Inches(0.8 + i * 2.98)
        y = y_start + Inches(0.1)
        w = Inches(2.8)
        h = Inches(4.8)
        
        c = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        c.fill.solid()
        c.fill.fore_color.rgb = CARD_BG
        c.line.color.rgb = col
        c.line.width = Pt(1.5)
        
        tb = slide7.shapes.add_textbox(x + Inches(0.15), y + Inches(0.15), w - Inches(0.3), h - Inches(0.3))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p1 = tf.paragraphs[0]
        p1.text = p_tag
        p1.font.size = Pt(10)
        p1.font.bold = True
        p1.font.color.rgb = col
        
        p2 = tf.add_paragraph()
        p2.text = p_title
        p2.font.size = Pt(13)
        p2.font.bold = True
        p2.font.color.rgb = TEXT_WHITE
        
        p3 = tf.add_paragraph()
        p3.text = "\nKey Deliverables:"
        p3.font.size = Pt(10)
        p3.font.bold = True
        p3.font.color.rgb = TEXT_MUTED
        
        for b in p_bullets:
            pb = tf.add_paragraph()
            pb.text = f"• {b}"
            pb.font.size = Pt(9.5)
            pb.font.color.rgb = TEXT_SUB
            
        ps = tf.add_paragraph()
        ps.text = f"\nStatus: {status}"
        ps.font.size = Pt(10)
        ps.font.bold = True
        ps.font.color.rgb = col

    # =============================================================
    # SLIDE 8: Expected Impact & Quantifiable Business ROI
    # =============================================================
    slide8 = prs.slides.add_slide(blank_slide_layout)
    set_bg(slide8)
    y_start = add_slide_header(
        slide8,
        "7. Expected Impact & Quantifiable Business ROI",
        "Massive Value Creation for Engineering & Talent Teams",
        "Transforming candidate experience while returning thousands of productive engineering hours to the business.",
        AMBER_GOLD
    )
    
    # 4 Huge Metric Cards
    metrics_4 = [
        ("80%", "Engineering Time Saved", "Frees senior architects from repetitive initial screening loops.", AGORA_BLUE),
        ("3.8x", "Faster Time-to-Offer", "Eliminates scheduling lag; candidates interview 24/7 on-demand.", CYAN_ACCENT),
        ("$1.2M+", "Annual Cost Savings", "Calculated on 100 technical hires @ $200/hr engineer rate.", EMERALD_GREEN),
        ("100%", "Objective & Auditable", "Every decision grounded in verbatim transcript evidence.", PURPLE_ACCENT)
    ]
    
    for i, (stat, label, desc, col) in enumerate(metrics_4):
        x = Inches(0.8 + i * 2.98)
        y = y_start + Inches(0.1)
        w = Inches(2.8)
        h = Inches(2.4)
        
        c = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        c.fill.solid()
        c.fill.fore_color.rgb = CARD_BG
        c.line.color.rgb = col
        c.line.width = Pt(1.5)
        
        tb = slide8.shapes.add_textbox(x + Inches(0.12), y + Inches(0.12), w - Inches(0.24), h - Inches(0.24))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p1 = tf.paragraphs[0]
        p1.text = stat
        p1.font.size = Pt(36)
        p1.font.bold = True
        p1.font.color.rgb = col
        p1.alignment = PP_ALIGN.CENTER
        
        p2 = tf.add_paragraph()
        p2.text = label
        p2.font.size = Pt(12)
        p2.font.bold = True
        p2.font.color.rgb = TEXT_WHITE
        p2.alignment = PP_ALIGN.CENTER
        
        p3 = tf.add_paragraph()
        p3.text = desc
        p3.font.size = Pt(9.5)
        p3.font.color.rgb = TEXT_MUTED
        p3.alignment = PP_ALIGN.CENTER

    # Candidate vs Enterprise Row
    cands_box = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.5), Inches(5.733), Inches(1.8))
    cands_box.fill.solid()
    cands_box.fill.fore_color.rgb = CARD_ALT
    cands_box.line.color.rgb = AGORA_BLUE
    
    tb_cand = slide8.shapes.add_textbox(Inches(0.95), Inches(4.6), Inches(5.433), Inches(1.6))
    tf_c = tb_cand.text_frame
    tf_c.word_wrap = True
    p_c1 = tf_c.paragraphs[0]
    p_c1.text = "FOR CANDIDATES"
    p_c1.font.size = Pt(11)
    p_c1.font.bold = True
    p_c1.font.color.rgb = AGORA_BLUE
    
    p_c2 = tf_c.add_paragraph()
    p_c2.text = "Zero scheduling friction, instant fair feedback, realistic interview pressure without human bias, and an actionable development roadmap regardless of outcome."
    p_c2.font.size = Pt(11)
    p_c2.font.color.rgb = TEXT_SUB

    org_box = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(4.5), Inches(5.733), Inches(1.8))
    org_box.fill.solid()
    org_box.fill.fore_color.rgb = CARD_ALT
    org_box.line.color.rgb = CYAN_ACCENT
    
    tb_org = slide8.shapes.add_textbox(Inches(6.95), Inches(4.6), Inches(5.433), Inches(1.6))
    tf_o = tb_org.text_frame
    tf_o.word_wrap = True
    p_o1 = tf_o.paragraphs[0]
    p_o1.text = "FOR HIRING ORGANIZATIONS"
    p_o1.font.size = Pt(11)
    p_o1.font.bold = True
    p_o1.font.color.rgb = CYAN_ACCENT
    
    p_o2 = tf_o.add_paragraph()
    p_o2.text = "Standardized calibration across global locations, immediate deep dive into candidates' real technical depth, and elimination of false-positive senior hires."
    p_o2.font.size = Pt(11)
    p_o2.font.color.rgb = TEXT_SUB

    # =============================================================
    # SLIDE 9: Conclusion, Live Demo & Q&A
    # =============================================================
    slide9 = prs.slides.add_slide(blank_slide_layout)
    set_bg(slide9)
    y_start = add_slide_header(
        slide9,
        "✨ Summary & Live Demo Ready",
        "Pioneering the Future of Real-Time Voice Conversational AI",
        "Vocalis AI proves what is possible when multi-agent intelligence meets Agora's world-class ultra-low latency real-time voice infrastructure.",
        EMERALD_GREEN
    )
    
    # 5-Box Evaluation Criteria Grid
    eval_boxes = [
        ("INNOVATION", "Multi-Agent", "Deliberative Turn-Taking", AGORA_BLUE),
        ("RELEVANCE", "High ROI", "$400B Hiring Market", CYAN_ACCENT),
        ("FEASIBILITY", "Working Prototype", "Fully Built & Tested", INDIGO_BLUE),
        ("IMPACT", "80% Time Saved", "$1.2M Saved / 100 Hires", EMERALD_GREEN),
        ("AGORA STACK", "Native SD-RTN", "<100ms Voice + VAD", AMBER_GOLD)
    ]
    
    for i, (tag, title, desc, col) in enumerate(eval_boxes):
        x = Inches(0.8 + i * 2.38)
        y = y_start + Inches(0.2)
        w = Inches(2.2)
        h = Inches(2.2)
        
        c = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        c.fill.solid()
        c.fill.fore_color.rgb = CARD_BG
        c.line.color.rgb = col
        c.line.width = Pt(1.5)
        
        tb = slide9.shapes.add_textbox(x + Inches(0.1), y + Inches(0.15), w - Inches(0.2), h - Inches(0.3))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p1 = tf.paragraphs[0]
        p1.text = tag
        p1.font.size = Pt(10)
        p1.font.bold = True
        p1.font.color.rgb = col
        p1.alignment = PP_ALIGN.CENTER
        
        p2 = tf.add_paragraph()
        p2.text = title
        p2.font.size = Pt(13)
        p2.font.bold = True
        p2.font.color.rgb = TEXT_WHITE
        p2.alignment = PP_ALIGN.CENTER
        
        p3 = tf.add_paragraph()
        p3.text = desc
        p3.font.size = Pt(9.5)
        p3.font.color.rgb = TEXT_MUTED
        p3.alignment = PP_ALIGN.CENTER

    # Action Banner at Bottom
    action_card = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.5), Inches(11.733), Inches(1.8))
    action_card.fill.solid()
    action_card.fill.fore_color.rgb = CARD_BG
    action_card.line.color.rgb = EMERALD_GREEN
    action_card.line.width = Pt(1.5)
    
    tb_act = slide9.shapes.add_textbox(Inches(1.0), Inches(4.7), Inches(11.333), Inches(1.4))
    tf_a = tb_act.text_frame
    tf_a.word_wrap = True
    
    pa1 = tf_a.paragraphs[0]
    pa1.text = "🎙️ Launch Interactive Voice Studio & Live Evaluation Deck"
    pa1.font.size = Pt(18)
    pa1.font.bold = True
    pa1.font.color.rgb = CYAN_ACCENT
    pa1.alignment = PP_ALIGN.CENTER
    
    pa2 = tf_a.add_paragraph()
    pa2.text = "\nThank you for reviewing our submission to the Echosphere Hackathon."
    pa2.font.size = Pt(13)
    pa2.font.color.rgb = TEXT_WHITE
    pa2.alignment = PP_ALIGN.CENTER
    
    pa3 = tf_a.add_paragraph()
    pa3.text = "We welcome questions from the evaluation committee!"
    pa3.font.size = Pt(12)
    pa3.font.bold = True
    pa3.font.color.rgb = EMERALD_GREEN
    pa3.alignment = PP_ALIGN.CENTER
    
    output_path = r"d:\Riyanshi\01_coding\projects\37 VoiceIntro AI\presentation.pptx"
    try:
        prs.save(output_path)
        print(f"Direct HTML to PPTX Conversion successful: {output_path}")
    except PermissionError:
        output_path_alt = r"d:\Riyanshi\01_coding\projects\37 VoiceIntro AI\Vocalis_AI_Echosphere_Presentation.pptx"
        prs.save(output_path_alt)
        print(f"Saved to alternative path due to file lock: {output_path_alt}")


if __name__ == "__main__":
    create_deck_from_html()
